"""Dump all text content from each slide in the PPTX for analysis."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

prs = Presentation("Lecture2_intro.pptx")

for slide_idx, slide in enumerate(prs.slides):
    print(f"=== SLIDE {slide_idx + 1} === Layout: {slide.slide_layout.name}")
    for shape_idx, shape in enumerate(slide.shapes):
        stype = str(shape.shape_type)
        print(f"  Shape {shape_idx}: type={stype}, name={repr(shape.name)}")
        print(f"    pos=({shape.left}, {shape.top}), size=({shape.width}, {shape.height})")
        if shape.has_text_frame:
            tf = shape.text_frame
            for p_idx, para in enumerate(tf.paragraphs):
                text = para.text
                if not text.strip():
                    continue
                align = para.alignment
                bidi = None
                # Check for RTL
                pPr = para._pPr
                if pPr is not None:
                    bidi_attr = pPr.attrib.get('{http://schemas.openxmlformats.org/drawingml/2006/main}rtl', 
                                pPr.attrib.get('rtl', None))
                    if bidi_attr:
                        bidi = bidi_attr
                
                print(f"    P{p_idx}: align={align} bidi={bidi}")
                for r in para.runs:
                    font = r.font
                    bold = font.bold
                    size = font.size
                    color = None
                    try:
                        if font.color and font.color.type:
                            color = str(font.color.rgb)
                    except:
                        pass
                    lang = None
                    rPr = r._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                    if rPr is not None:
                        lang = rPr.attrib.get('lang', None)
                    t = r.text
                    print(f"      Run: text={repr(t[:80])} bold={bold} sz={size} color={color} lang={lang}")
        else:
            # Check if it's an image
            if shape.shape_type == 13:  # Picture
                print(f"    [IMAGE: {shape.image.content_type}]")
            elif hasattr(shape, 'image'):
                try:
                    print(f"    [IMAGE: {shape.image.content_type}]")
                except:
                    pass
    print()
