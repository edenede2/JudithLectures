"""
translate_pptx.py — Translate a Hebrew PPTX to English.

Usage:
    python translate_pptx.py [input.pptx] [output.pptx]
    
Defaults:
    input:  Lecture2_intro.pptx
    output: Lecture2_intro_EN.pptx

Strategy:
  1. Open the original PPTX (all images, shapes, layouts preserved).
  2. For each slide in the translation map, find the target shape.
  3. Replace paragraph text while PRESERVING original run formatting
     (font size, color, bold, etc.) from the first run of each paragraph.
  4. Convert RTL alignment → LTR, he-IL lang → en-US.
  5. Save as a new file.
"""

import sys
import copy
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree


# ─────────────────────────────────────────────────────────────
# Translation map
# ─────────────────────────────────────────────────────────────
# slide_num -> [ (shape_name, [paragraph_specs, ...]), ... ]
#
# paragraph_spec is one of:
#   None                → keep the original paragraph unchanged
#   "CLEAR"             → clear the paragraph text (make empty)
#   "text"              → single run, inherit formatting from original
#   ("text",)           → same
#   [("text", {}), ...] → multiple runs with optional overrides:
#                          bold=True/False, color="FF0000"
#
# The number of paragraph_specs should match the original paragraph
# count in that shape.  Extra original paragraphs are cleared.
# ─────────────────────────────────────────────────────────────


def build_translations():
    """Return the translation map."""
    T = {}

    # ── SLIDE 1 ── Title slide
    T[1] = [
        ("Title 1", [
            "Project Life Cycle in Data Science",
        ]),
        ("Subtitle 2", [
            "Judith Somekh",
        ]),
    ]

    # Slides 2, 3: already English / diagram only

    # ── SLIDE 4 ── Section divider
    T[4] = [
        ("Title 1", [
            [("2. Data Gathering  & Understanding\nData collection & understanding", {})],
        ]),
    ]

    # ── SLIDE 5 ──
    T[5] = [
        ("Content Placeholder 2", [
            "The data include the raw material available from which the solution will be built.",
            "Usually, finding the right data takes time and effort.",
            "Data are often collected from multiple sources.",
            "CLEAR",  # empty line
            "What are some of the problems that can exist in data?",
        ]),
    ]

    # ── SLIDE 6 ──
    T[6] = [
        ("Content Placeholder 2", [
            "Some of the questions worth considering are:",
            "What data do I need for my project?",
            "Where are the data located?",
            "How can I obtain them?",
            "What are the costs?",
            "What is the most efficient way to store and access all of this?",
            "How do we preserve data privacy and confidentiality?",
        ]),
    ]

    # ── SLIDE 7 ──
    T[7] = [
        ("Content Placeholder 2", [
            "It is important to understand the strengths and limitations of the data, because there is rarely an exact match between the data and the problem.",
            "Historical data are often collected for purposes unrelated to the current business problem, or with no explicit purpose at all (for example, medical records).",
            "For example, a customer database, a transaction database, and a marketing response database contain different information, may cover different overlapping populations, and may have different levels of reliability.",
            "Data costs may vary. Some data will be available for free, while others will need to be purchased or will require effort to obtain (for example, GEO, the Israeli Digital Health Project).",
            "Sometimes the data do not exist, and entire projects will be required to organize their collection (for example, gene expression data require performing experiments).",
            "A critical part of the data understanding stage is estimating the costs and benefits of each data source.",
        ]),
    ]

    # ── SLIDE 8 ──
    T[8] = [
        ("Title 1", [
            "Free medical information for researchers and start-up companies",
        ]),
    ]

    # ── SLIDE 9 ──
    T[9] = [
        ("Content Placeholder 2", [
            "As data understanding progresses, solution paths may change direction in response, and the team\u2019s efforts may even split.",
            "Example: Fraud detection \u2013 data mining is widely used for fraud detection, and many fraud detection problems include classic data mining tasks.",
            "Sometimes fraud problems that initially look similar turn out to be fundamentally different during the data understanding process.",
            [("Catching credit card fraud:", {"bold": True})],
            "How would we do that?",
        ]),
    ]

    # ── SLIDE 10 ──
    T[10] = [
        ("Content Placeholder 2", [
            "Catching credit card fraud:",
            "To build a predictive model for fraud, we need to label actions as fraudulent / legitimate.",
            "Question: Who will perform the labeling?",
            "How will we identify fraud? Who will identify it?",
            "Is the identification reliable?",
        ]),
    ]

    # ── SLIDE 11 ──
    T[11] = [
        ("Content Placeholder 2", [
            "Catching credit card fraud:",
            "Charges appear in each customer\u2019s account, so fraudulent charges are usually detected\u2014if not initially by the company, then later by the customer when reviewing account activity.",
            "We can assume that almost all fraud is detected and labeled reliably, because the legitimate customer and the person committing the fraud are different people with opposite interests.",
            "Credit card transactions have reliable labels (fraud or legitimate) that can be used to train a machine learning model.",
        ]),
    ]

    # ── SLIDE 12 ──
    T[12] = [
        ("Content Placeholder 2", [
            "CLEAR",
            "CLEAR",
            "Medicare fraud: this is a huge problem in the United States, costing billions of dollars each year.",
            "The fraudsters are healthcare providers who submit false claims for various treatments for patients.",
            "CLEAR",
            "Is this a conventional fraud detection problem?",
        ]),
    ]

    # ── SLIDE 13 ──
    T[13] = [
        ("Content Placeholder 2", [
            "CLEAR",
            "CLEAR",
            "Medicare fraud: this is a huge problem in the United States, costing billions of dollars each year.",
            "The fraudsters are healthcare providers who submit false claims for various treatments for patients.",
            "CLEAR",
            "Is this a conventional fraud detection problem?",
            "When we examine the relationship between the business problem and the data, we understand that the problem is significantly different.",
            "Why? How will we identify fraud? Who identifies it?",
        ]),
    ]

    # ── SLIDE 14 ──
    T[14] = [
        ("Content Placeholder 2", [
            "CLEAR",
            "CLEAR",
            "Example: Medicare fraud",
            "The perpetrators of fraud\u2014healthcare providers submitting false claims, and sometimes their patients\u2014are also legitimate service providers using the billing system.",
            [("The people committing the fraud are a subset of legitimate users; there is no separate disinterested party that will accurately declare which claims are \u201ccorrect.\u201d", {"bold": True})],
            "As a result, Medicare billing data do not contain a reliable target variable indicating fraud, so a supervised learning approach that might work for credit card fraud is not relevant here.",
            [("Such a problem usually requires ", {}),
             ("unsupervised", {"bold": True}),
             (" approaches such as clustering, anomaly detection, and co-occurrence grouping of similar data.", {})],
        ]),
    ]

    # ── SLIDE 15 ──
    T[15] = [
        ("Content Placeholder 2", [
            "Statistical learning lies at the intersection of computer science and statistics and refers to a set of tools for understanding data. It combines statistical methods, computational methods, and methods from machine learning.",
            "These tools are categorized as:",
            [("Supervised learning", {"bold": True})],
            "Building a statistical model for predicting or estimating an output based on one or more inputs.",
            [("Unsupervised learning", {"bold": True})],
            "Inputs exist, but there is no known corresponding output. Even so, we can learn about relationships and structure from the data. There is no response variable that \u201csupervises\u201d the analysis.",
        ]),
    ]

    # ── SLIDE 16 ──
    T[16] = [
        ("Content Placeholder 2", [
            [("The fact that both cases are called fraud detection problems is based on a ", {}),
             ("superficial similarity that is actually misleading", {"bold": True}),
             (", because they are two different cases requiring ", {}),
             ("different analyses", {"bold": True}),
             (".", {})],
            [("In data understanding, we must dig beneath the surface ", {"bold": True}),
             ("to uncover the structure of the business problem and the available data for understanding it, and then match them to one or more data mining and analysis tasks requiring the knowledge, science, and technology we possess.", {})],
            "It is not unusual for a business problem to contain several data mining and analysis tasks, often of different kinds, requiring integration of the solutions to the different tasks.",
        ]),
    ]

    # Slide 17: diagram only

    # ── SLIDE 18 ── Section divider
    T[18] = [
        ("Title 1", [
            [("3. Data Cleaning & Preparation\nCleaning and preparing data", {})],
        ]),
    ]

    # ── SLIDE 19 ──
    T[19] = [
        ("Content Placeholder 2", [
            "Powerful data analysis techniques exist, but they are still limited by certain requirements imposed on the data they use.",
            "Data entered into analysis usually need to be represented differently from the original form in which they were collected or delivered, and this often requires conversion and cleaning efforts.",
            "What kinds of problems can our data have?",
        ]),
    ]

    # Slides 20, 21: already English

    # ── SLIDE 22 ──
    T[22] = [
        ("Title 1", [
            "A Mars lander crashed because of a unit conversion error",
        ]),
    ]

    # ── SLIDE 23 ──
    T[23] = [
        ("Title 1", [
            "Data Leakage",
        ]),
        ("Content Placeholder 2", [
            "Leakage is a situation in which a variable collected in historical data provides information about the target variable\u2014information that appears in the historical data but is not actually available when a decision must be made.",
            "Example 1: We want to predict whether a customer will be a \u201cbig spender\u201d; knowing the categories of items purchased (or the amount of tax paid) may be highly predictive, but these are not known at decision time (Kohavi & Parekh, 2003).",
            "Example 2: When predicting whether, at a certain point in time, a website visitor will end her session or continue to another page, the variable \u201ctotal number of webpages visited in the session\u201d is predictive. However, that total is only known after the session ends (Kohavi et al., 2000)\u2014at which point the target is already known.",
        ]),
    ]

    # Slide 24: diagram only

    # ── SLIDE 25 ── Section divider
    T[25] = [
        ("Title 1", [
            [("4. Data Exploration & Modeling ", {"bold": True}),
             ("\nData exploration and models", {})],
        ]),
    ]

    # ── SLIDE 26 ──
    T[26] = [
        ("Content Placeholder 3", [
            "The actual data analysis process is based on computational techniques, statistical methods, and algorithms.",
            "The result of modeling is some form of model that finds patterns and structures in the data.",
        ]),
    ]

    # ── SLIDE 27 ──
    T[27] = [
        ("Content Placeholder 3", [
            [("We must ", {}),
             ("understand the patterns, problems, and biases in the data", {"bold": True}),
             (":", {})],
            "This is done by sampling and analyzing a random subset of the data from the full dataset.",
            "Use histograms or distribution curves to see the general trend, or even build an interactive visualization that allows drilling down to individual data points.",
            "Investigate the story behind outliers.",
            [("Begin forming hypotheses", {"bold": True}),
             (" about the data and the problem we are trying to solve or analyze.", {})],
            "Example: In predicting students\u2019 grades, you could examine the relationship between grades and sleep.",
            "Example: In predicting real-estate prices, you could draw prices as a spatial heatmap to see whether you can detect trends.",
        ]),
    ]

    # Slide 28: already English

    # ── SLIDE 29 ──
    T[29] = [
        ("Content Placeholder 2", [
            "Patients with cancer and healthy individuals",
            "The features of each patient are represented by thousands of genes (measured mRNA levels), and the levels change between healthy and diseased states",
        ]),
    ]

    # ── SLIDE 30 ──
    T[30] = [
        ("Title 1", [
            [("4.2 Feature Engineering", {"bold": True})],
        ]),
        ("Content Placeholder 3", [
            "In machine learning, a feature is a measurable property or characteristic of an observed phenomenon.",
            "For predicting student grades, one possible feature is the amount of sleep at night.",
            "Question: what other relevant feature ideas can you think of?",
            "In more complex tasks such as face recognition, features can be histograms counting the number of black pixels.",
            "According to Andrew Ng: \u201cComing up with features is difficult, time-consuming, requires expert knowledge. \u2018Applied machine learning\u2019 is basically feature engineering.\u201d",
        ]),
    ]

    # ── SLIDE 31 ──
    T[31] = [
        ("Title 1", [
            [("4.2 Feature Engineering", {"bold": True})],
        ]),
        ("Content Placeholder 3", [
            "Feature engineering is the process of using domain knowledge to transform raw data into informative features that represent the business problem being solved.",
            "This stage directly affects the accuracy of the predictive model built in the next stage.",
            "Usually we perform two types of tasks in feature engineering:",
            "Feature selection \u2013 choosing features relevant to the analysis",
            "Feature construction \u2013 for example, combining features into one feature",
        ]),
    ]

    # ── SLIDE 32 ──
    T[32] = [
        ("Title 1", [
            [("4.2 Feature Selection", {"bold": True})],
        ]),
        ("Content Placeholder 3", [
            "Feature selection is the process of removing features that add more noise than information.",
            "For this purpose we use filtering methods:",
            "Filter methods \u2013 generate a statistical measure to assign a score to each feature.",
            [("Wrapper methods", {}),
             (" \u2013 treat feature selection as a search problem and use heuristics to perform the search.", {})],
            [("Embedded methods", {}),
             (" \u2013 use machine learning itself to determine which features contribute best to accuracy.", {})],
        ]),
    ]

    # ── SLIDE 33 ──
    T[33] = [
        ("Title 1", [
            [("4.2 Feature Construction", {"bold": True})],
        ]),
        ("Content Placeholder 3", [
            "Creating new features from original features.",
            "For example, when you have a continuous variable but domain knowledge says that only crossing a certain threshold matters.",
            "Given feature: subject\u2019s age.",
            "If the model only cares whether a person is an adult or a minor, we can place a threshold at age 18 and assign different categories to cases above and below it.",
            "Other threshold examples: sick/healthy, obese/thin.",
            "Merging multiple features to make them more informative by taking their sum, difference, or product.",
            "Example: predicting student grades.",
            "Given: features representing number of hours of sleep for a student each night.",
            "What processing would you perform on these features for prediction?",
        ]),
    ]

    # ── SLIDE 34 ──
    T[34] = [
        ("Title 1", [
            [("4.2 Feature Construction", {"bold": True})],
        ]),
        ("Content Placeholder 3", [
            "Creating new features from original features.",
            "For example, when you have a continuous variable but domain knowledge says that only crossing a certain threshold matters.",
            "Given feature: subject\u2019s age.",
            "If the model only cares whether a person is an adult or a minor, we can place a threshold at age 18 and assign different categories to cases above and below it.",
            "Other threshold examples: sick/healthy, obese/thin.",
            "Merging multiple features to make them more informative by taking their sum, difference, or product.",
            "Example: predicting student grades.",
            "Given: features representing number of hours of sleep for a student each night.",
            "What processing would you perform on these features for prediction?",
            "You should create a feature representing the student\u2019s average amount of sleep.",
        ]),
    ]

    # ── SLIDE 35 ──
    T[35] = [
        ("Title 1", [
            [("4.3 Predictive Modeling", {"bold": True})],
        ]),
        ("Content Placeholder 3", [
            "The actual analysis work!",
            "Includes:",
            "Machine learning techniques for building a predictive model and testing its accuracy.",
            "Use of comprehensive statistical methods and tests to ensure that the model\u2019s results are sensible and meaningful.",
            [("Based on the questions asked during the business understanding stage, this is where we decide ", {}),
             ("which model to choose for the problem.", {"bold": True})],
            "This is never an easy decision, and there is no single correct answer.",
            "The chosen model (or models\u2014and it is always preferable to test several) will depend on the size, type, and quality of your data, how much time and computing resources you are willing to invest, and the type of output you intend to produce.",
        ]),
    ]

    # ── SLIDE 36 ──
    T[36] = [
        ("Content Placeholder 3", [
            "After training your model, it is critical to evaluate its quality and its success in prediction.",
            "Cross-validation\u2014usually k-fold cross-validation\u2014is commonly used to measure model accuracy.",
            "This involves splitting the dataset into equally sized groups of instances, building a model on all but one group, and repeating the process with different held-out groups. This allows the model to be trained across all the data rather than relying on a simple train-test split.",
            "Charts such as ROC curves, which show the true positive rate against the false positive rate, are also used to measure model success.",
        ]),
    ]

    # ── SLIDE 37 ──
    T[37] = [
        ("Title 1", [
            [("4.4 ", {"bold": True}),
             ("Visualization", {})],
        ]),
        ("Content Placeholder 3", [
            "Data visualization seems simple, but often it is one of the most challenging tasks.",
            "Data visualization combines communication, psychology, statistics, and art.",
            "The ultimate goal is to communicate the data in a simple, effective, and visually pleasing way.",
            "After extracting the intended insights from your model, you must represent them in a way that the various stakeholders in the project can understand.",
            "The chart should be clear and tell the story!",
        ]),
    ]

    # Slide 38: diagram only

    # ── SLIDE 39 ──
    T[39] = [
        ("TextBox 6", [
            "Goal 1: Rigorously evaluate the results of the data analysis and ensure that they are valid and reliable before moving on.",
            "If we search enough, we will find patterns in any dataset\u2014but they may not survive more careful testing.",
            [("We want to be sure that the models and patterns extracted from the data are real regularities and not just strange phenomena or sampling artifacts.", {})],
        ]),
    ]

    # ── SLIDE 40 ──
    T[40] = [
        ("TextBox 6", [
            "CLEAR",
            [("Controlled evaluation", {"bold": True}),
             (": first test a model within a controlled \u201claboratory\u201d process.", {})],
            [("Evaluation of external considerations", {"bold": True}),
             (": even if a model passes strict evaluation tests in a controlled environment, there may still be external considerations that make it impractical.", {})],
            "For example, a common flaw in detection analyses (such as fraud detection, spam detection, and intrusion monitoring) is that they generate too many false alarms.",
            "Additional considerations must be addressed, such as: how much it will cost the team to deal with all those false alarms, and what the meaning and cost of customer dissatisfaction will be.",
        ]),
    ]

    # ── SLIDE 41 ──
    T[41] = [
        ("TextBox 6", [
            [("Evaluation in the real environment", {"bold": True}),
             (": in some cases we will extend evaluation into the real environment, for example by instrumenting a live system to allow randomized experiments.", {})],
            "CLEAR",
            "CLEAR",
            "Why do you think this is necessary?",
        ]),
    ]

    # ── SLIDE 42 ──
    T[42] = [
        ("TextBox 6", [
            [("Evaluation in the real environment", {"bold": True}),
             (": in some cases we will extend evaluation into the real environment, for example by using instrumentation in a live system that enables randomized experiments.", {})],
            "In the example of customer churn in cellular companies\u2014if laboratory testing suggests that a data-driven model can help reduce churn, we may want to move to real-world evaluation in a live system that applies the model randomly to some customers while keeping others as a control group. Such experiments must be carefully planned.",
            [("Considering the changing environment", {"bold": True}),
             (":", {})],
            "What might change when the model is implemented in the company?",
        ]),
    ]

    # ── SLIDE 43 ──
    T[43] = [
        ("TextBox 6", [
            [("Evaluation in the real environment", {"bold": True}),
             (": in some cases we will extend evaluation into the real environment, for example by using instrumentation in a live system that enables randomized experiments.", {})],
            "In the example of customer churn in cellular companies\u2014if laboratory testing suggests that a data-driven model can help reduce churn, we may want to move to real-world evaluation in a live system that applies the model randomly to some customers while keeping others as a control group.",
            [("Considering the changing environment", {"bold": True}),
             (": We may also want to train our systems to ensure that the world does not change in a way that harms the model\u2019s decision-making. For example, behavior can change\u2014in some cases, such as ", {}),
             ("fraud or spam", {"bold": True}),
             (", directly in response to deployment of the models.", {})],
        ]),
    ]

    # ── SLIDE 44 ──
    T[44] = [
        ("Content Placeholder 2", [
            [("We want to ensure that the model meets the original business goals.", {"bold": True})],
            "After going through the entire lifecycle, it is time to return to the drawing board. Remember, this is a circular lifecycle, so it is an iterative process. This is the point where you evaluate how the success of your model relates back to your original understanding of the business and its goals.",
            "Does it address the problems that were identified?",
            "Does the analysis yield solid solutions?",
            [("If new insights arose during the first iteration of the lifecycle (which usually happens), you can now ", {}),
             ("inject that new knowledge", {"bold": True}),
             (" into the next iteration in order to produce even stronger insights and use the power of data to generate strong and exceptional outcomes for your business or project.", {})],
        ]),
    ]

    # ── SLIDE 45 ──
    T[45] = [
        ("Title 1", [
            [("Q: What do data scientists spend the most time doing?", {"bold": True}),
             ("\nWhat percentage of the time do you think each task takes?", {})],
        ]),
    ]

    # Slides 46, 47: already English

    return T


# ─────────────────────────────────────────────────────────────
# Core translation engine
# ─────────────────────────────────────────────────────────────

def _clone_rpr(source_rPr):
    """Deep-copy a <a:rPr> element, stripping he-IL lang."""
    if source_rPr is None:
        rPr = etree.SubElement(etree.Element('dummy'), qn('a:rPr'))
        rPr.set('lang', 'en-US')
        rPr.set('dirty', '0')
        return rPr
    new_rPr = copy.deepcopy(source_rPr)
    # Fix language
    if new_rPr.get('lang', '') == 'he-IL':
        new_rPr.set('lang', 'en-US')
    if new_rPr.get('altLang', '') == 'he-IL':
        new_rPr.set('altLang', 'en-US')
    return new_rPr


def _get_template_rpr(para):
    """Get a template <a:rPr> from the first run of a paragraph."""
    r_elems = para._p.findall(qn('a:r'))
    if r_elems:
        rPr = r_elems[0].find(qn('a:rPr'))
        return rPr
    return None


def _clear_runs(p_elem):
    """Remove all <a:r>, <a:br>, and <a:endParaRPr> from a <a:p> element."""
    for tag in (qn('a:r'), qn('a:br'), qn('a:endParaRPr')):
        for child in p_elem.findall(tag):
            p_elem.remove(child)


def _fix_ppr_ltr(p_elem):
    """Set paragraph properties to LTR."""
    pPr = p_elem.find(qn('a:pPr'))
    if pPr is not None:
        # Remove RTL attribute
        for attr in list(pPr.attrib.keys()):
            if 'rtl' in attr.lower():
                del pPr.attrib[attr]


def _add_run(p_elem, text, template_rPr, overrides=None):
    """
    Add a <a:r> element to p_elem.
    
    template_rPr: a source <a:rPr> element to clone formatting from.
    overrides: dict with optional keys 'bold', 'color' to override.
    """
    r_elem = etree.SubElement(p_elem, qn('a:r'))
    
    # Clone formatting from template
    new_rPr = _clone_rpr(template_rPr)
    # Ensure lang is en-US
    new_rPr.set('lang', 'en-US')
    new_rPr.set('dirty', '0')
    
    # Apply overrides
    if overrides:
        if 'bold' in overrides:
            if overrides['bold']:
                new_rPr.set('b', '1')
            else:
                # Remove bold if explicitly set to False
                if 'b' in new_rPr.attrib:
                    del new_rPr.attrib['b']
        if 'color' in overrides and overrides['color']:
            # Remove existing solidFill
            for sf in new_rPr.findall(qn('a:solidFill')):
                new_rPr.remove(sf)
            solidFill = etree.SubElement(new_rPr, qn('a:solidFill'))
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgbClr.set('val', overrides['color'])
    
    r_elem.insert(0, new_rPr)
    
    t_elem = etree.SubElement(r_elem, qn('a:t'))
    t_elem.text = text
    if text and (text[0] == ' ' or text[-1] == ' ' or '  ' in text):
        t_elem.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')


def _set_alignment(p_elem, align_str):
    """Set paragraph alignment."""
    align_map = {"LEFT": "l", "CENTER": "ctr", "RIGHT": "r"}
    pPr = p_elem.find(qn('a:pPr'))
    if pPr is None:
        pPr = etree.SubElement(p_elem, qn('a:pPr'))
        p_elem.insert(0, pPr)
    if align_str in align_map:
        pPr.set('algn', align_map[align_str])


def replace_shape_text(shape, para_specs):
    """
    Replace text in a shape while preserving original formatting.
    
    para_specs: list of paragraph specifications (see build_translations).
    """
    if not shape.has_text_frame:
        return

    txBody = shape.text_frame._txBody
    orig_p_elems = txBody.findall(qn('a:p'))

    # Gather template formatting from existing paragraphs
    templates = []
    for p in orig_p_elems:
        rPr = None
        for r in p.findall(qn('a:r')):
            rPr = r.find(qn('a:rPr'))
            if rPr is not None:
                break
        templates.append(rPr)

    # Determine target paragraph count
    target_count = len(para_specs)
    current_count = len(orig_p_elems)

    # Add paragraphs if needed (clone last one's structure)
    while current_count < target_count:
        new_p = copy.deepcopy(orig_p_elems[-1]) if orig_p_elems else etree.SubElement(txBody, qn('a:p'))
        _clear_runs(new_p)
        txBody.append(new_p)
        # Use the last known template
        templates.append(templates[-1] if templates else None)
        orig_p_elems = txBody.findall(qn('a:p'))
        current_count += 1

    # Remove excess paragraphs
    orig_p_elems = txBody.findall(qn('a:p'))
    while current_count > target_count and target_count > 0:
        txBody.remove(orig_p_elems[current_count - 1])
        orig_p_elems = txBody.findall(qn('a:p'))
        templates.pop()
        current_count -= 1

    # Process each paragraph
    orig_p_elems = txBody.findall(qn('a:p'))
    for i, spec in enumerate(para_specs):
        p_elem = orig_p_elems[i]
        template_rPr = templates[i] if i < len(templates) else templates[-1] if templates else None

        # Fix RTL → LTR
        _fix_ppr_ltr(p_elem)

        if spec is None:
            # Keep original paragraph unchanged
            continue

        # Clear existing runs
        _clear_runs(p_elem)

        if spec == "CLEAR":
            # Empty paragraph
            continue

        if isinstance(spec, str):
            # Single run, inherit all formatting
            _add_run(p_elem, spec, template_rPr)

        elif isinstance(spec, list):
            # Multiple runs with potential overrides
            for run_spec in spec:
                if isinstance(run_spec, tuple) and len(run_spec) == 2:
                    text, overrides = run_spec
                    _add_run(p_elem, text, template_rPr, overrides)
                elif isinstance(run_spec, str):
                    _add_run(p_elem, run_spec, template_rPr)

    # Set alignment: figure out from original; if was RTL/RIGHT, switch to LEFT
    for p_elem in txBody.findall(qn('a:p')):
        pPr = p_elem.find(qn('a:pPr'))
        if pPr is not None:
            algn = pPr.get('algn', '')
            if algn == 'r':
                pPr.set('algn', 'l')

    # Ensure <a:endParaRPr> is at the end of each paragraph (after all runs)
    for p_elem in txBody.findall(qn('a:p')):
        endParaRPr = p_elem.find(qn('a:endParaRPr'))
        if endParaRPr is not None:
            p_elem.remove(endParaRPr)
            if endParaRPr.get('lang', '') == 'he-IL':
                endParaRPr.set('lang', 'en-US')
            p_elem.append(endParaRPr)
        else:
            endParaRPr = etree.SubElement(p_elem, qn('a:endParaRPr'))
            endParaRPr.set('lang', 'en-US')
            endParaRPr.set('dirty', '0')


def fix_all_rtl(prs):
    """Convert any remaining RTL settings to LTR across the whole presentation."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                pPr = para._pPr
                if pPr is not None:
                    for attr in list(pPr.attrib.keys()):
                        if 'rtl' in attr.lower():
                            del pPr.attrib[attr]
                for run in para.runs:
                    rPr = run._r.find(qn('a:rPr'))
                    if rPr is not None:
                        if rPr.get('lang', '') == 'he-IL':
                            rPr.set('lang', 'en-US')
                        if rPr.get('altLang', '') == 'he-IL':
                            rPr.set('altLang', 'en-US')


def translate_pptx(input_path, output_path):
    """Main translation function."""
    print(f"Loading: {input_path}")
    prs = Presentation(input_path)

    translations = build_translations()

    for slide_num, shape_list in sorted(translations.items()):
        slide_idx = slide_num - 1
        if slide_idx >= len(prs.slides):
            print(f"  WARNING: Slide {slide_num} does not exist")
            continue

        slide = prs.slides[slide_idx]
        print(f"  Translating slide {slide_num}...")

        for shape_name, para_specs in shape_list:
            target = None
            for shape in slide.shapes:
                if shape.name == shape_name:
                    target = shape
                    break
            if target is None:
                print(f"    WARNING: Shape '{shape_name}' not found on slide {slide_num}")
                continue

            replace_shape_text(target, para_specs)

    print("  Fixing remaining RTL/language settings...")
    fix_all_rtl(prs)

    print(f"Saving: {output_path}")
    prs.save(output_path)
    print("Done!")


if __name__ == "__main__":
    if len(sys.argv) == 3:
        translate_pptx(sys.argv[1], sys.argv[2])
    else:
        translate_pptx("Lecture2_intro.pptx", "Lecture2_intro_EN.pptx")
