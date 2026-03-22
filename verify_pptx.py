"""Verify the translated PPTX: check for Hebrew and font sizes."""
import re
from pptx import Presentation

prs = Presentation('Lecture2_intro_EN.pptx')
hebrew = re.compile(r'[\u0590-\u05FF]')
found = False
for i, slide in enumerate(prs.slides):
    for s in slide.shapes:
        if s.has_text_frame:
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    if hebrew.search(r.text):
                        print(f"Slide {i+1} [{s.name}]: {repr(r.text[:60])}")
                        found = True
if not found:
    print("No Hebrew text remaining.")

print()

# Check font sizes on key slides
for check_slide in [1, 5, 6, 7, 9, 14, 39, 44]:
    slide = prs.slides[check_slide - 1]
    print(f"--- Slide {check_slide} ---")
    for s in slide.shapes:
        if s.has_text_frame:
            for p in s.text_frame.paragraphs:
                txt = p.text.strip()
                if txt:
                    sizes = []
                    for r in p.runs:
                        sz = r.font.size
                        sizes.append(str(sz))
                    # Also show alignment
                    algn = p.alignment
                    print(f"  [{algn}] {s.name}: sizes={sizes} text={repr(txt[:80])}")
