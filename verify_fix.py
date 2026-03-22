from pptx import Presentation
from lxml import etree

prs = Presentation('Lecture2_intro_EN.pptx')
ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
problem_found = False
slides_to_check = [1, 4, 5, 8, 18, 25, 26, 45]

for sn in slides_to_check:
    slide = prs.slides[sn-1]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for pi, para in enumerate(shape.text_frame.paragraphs):
            p = para._p
            children = list(p)
            tags = [c.tag.split('}')[-1] for c in children]
            if 'endParaRPr' in tags and 'r' in tags:
                epr_idx = tags.index('endParaRPr')
                r_idx = tags.index('r')
                if epr_idx < r_idx:
                    print("BUG: Slide %d %s P%d: endParaRPr before run" % (sn, shape.name, pi))
                    problem_found = True

if not problem_found:
    print('OK: All endParaRPr elements are correctly positioned after runs.')

print()
for sn in slides_to_check:
    slide = prs.slides[sn-1]
    print('=== SLIDE %d ===' % sn)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for pi, para in enumerate(shape.text_frame.paragraphs):
            if para.text.strip():
                print('  %s P%d: %s' % (shape.name, pi, repr(para.text[:90])))
    print()
