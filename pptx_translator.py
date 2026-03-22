"""
pptx_translator.py — General-purpose PPTX translation engine.

Workflow:
  1. generate_template()      — Extract text from a PPTX into a template MD
  2. User translates via ChatGPT (using the generated prompt + template)
  3. translate_pptx_bytes()    — Apply the translated MD back to the original PPTX

Template format (generated):
  > original text              — one line per paragraph

Translated format (filled by user/ChatGPT):
  - translated text            — one line per paragraph
  - **bold text**              — bold formatting
  - normal **bold** normal     — mixed formatting
  -                            — empty/spacer paragraph
  \\n within a line            — line break inside a paragraph
"""

import re
import copy
import io

from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree


# ─── Constants ────────────────────────────────────────────────

LANG_CODES = {
    "Hebrew": "he-IL",
    "Arabic": "ar-SA",
    "English": "en-US",
    "French": "fr-FR",
    "Spanish": "es-ES",
    "German": "de-DE",
    "Russian": "ru-RU",
    "Chinese (Simplified)": "zh-CN",
    "Japanese": "ja-JP",
    "Korean": "ko-KR",
    "Portuguese": "pt-BR",
    "Italian": "it-IT",
    "Dutch": "nl-NL",
    "Turkish": "tr-TR",
    "Polish": "pl-PL",
}

_RTL_PREFIXES = ("he", "ar", "fa", "ur")


# ─── Language helpers ─────────────────────────────────────────

def _is_rtl_lang(lang_code):
    if not lang_code:
        return False
    return lang_code.lower().startswith(_RTL_PREFIXES)


def _detect_source_lang(prs):
    """Return the most common non-English language code in the presentation."""
    counts = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    rPr = run._r.find(qn("a:rPr"))
                    if rPr is not None:
                        lang = rPr.get("lang", "")
                        if lang and lang != "en-US":
                            counts[lang] = counts.get(lang, 0) + 1
    return max(counts, key=counts.get) if counts else None


# ─── Template generation ─────────────────────────────────────

def _generate_template_from_prs(prs, source_lang, target_lang):
    lines = []
    lines.append(f"# Translation Template\n\n")
    lines.append(f"Source language: {source_lang}\n")
    lines.append(f"Target language: {target_lang}\n\n")
    lines.append("---\n")

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        lines.append(f"\n## Slide {slide_num}\n\n")

        text_shapes = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if "Slide Number" in shape.name:
                continue
            full_text = shape.text_frame.text.strip()
            if not full_text:
                continue
            text_shapes.append(shape)

        if not text_shapes:
            lines.append("SKIP\n\n")
            lines.append("---\n")
            continue

        for shape in text_shapes:
            lines.append(f"### Shape: {shape.name}\n")
            for para in shape.text_frame.paragraphs:
                parts = []
                for run in para.runs:
                    if run.font.bold:
                        parts.append(f"**{run.text}**")
                    else:
                        parts.append(run.text)
                text = "".join(parts)
                if text.strip():
                    lines.append(f"> {text}\n")
                else:
                    lines.append(">\n")
            lines.append("\n")

        lines.append("---\n")

    return "".join(lines)


def generate_template(pptx_path, source_lang="Hebrew", target_lang="English"):
    """Generate a translation template MD from a PPTX file path."""
    prs = Presentation(pptx_path)
    return _generate_template_from_prs(prs, source_lang, target_lang)


def generate_template_from_bytes(pptx_bytes, source_lang="Hebrew", target_lang="English"):
    """Generate a translation template MD from PPTX bytes."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    return _generate_template_from_prs(prs, source_lang, target_lang)


# ─── MD parser ────────────────────────────────────────────────

def _parse_bold_text(text):
    """Parse a paragraph string with optional **bold** markers into a para spec.

    Returns:
      ""           — empty paragraph (spacer)
      "plain text" — single run, inherit original formatting
      [(text, {overrides}), ...] — multiple runs with formatting
    """
    if text is None or text.strip() == "":
        return ""

    # Literal \n → real newline
    text = text.replace("\\n", "\n")

    # Split on bold markers
    parts = re.split(r"(\*\*.+?\*\*)", text)

    if len(parts) == 1:
        return text

    runs = []
    for part in parts:
        if not part:
            continue
        if part.startswith("**") and part.endswith("**"):
            runs.append((part[2:-2], {"bold": True}))
        else:
            runs.append((part, {}))

    # Simplify: if only one non-bold run, return plain string
    if len(runs) == 1 and not runs[0][1].get("bold"):
        return runs[0][0]

    return runs


def parse_translation_md(md_text):
    """Parse a translated MD file.

    Supports both formats:
      - ``### Shape: Name`` and ``Shape: Name`` headers
      - ``- paragraph`` and ``Text: paragraph`` content lines

    Returns: {slide_num: [(shape_name, [para_specs])]}
    """
    result = {}
    current_slide = None
    current_shape = None
    current_paras = []

    def _save():
        nonlocal current_shape, current_paras
        if current_slide is not None and current_shape and current_paras:
            if current_slide not in result:
                result[current_slide] = []
            result[current_slide].append((current_shape, list(current_paras)))
        current_shape = None
        current_paras = []

    for line in md_text.split("\n"):
        line = line.rstrip()

        if not line or line.strip() == "---":
            continue

        # Slide header
        m = re.match(r"^##\s+Slide\s+(\d+)", line)
        if m:
            _save()
            current_slide = int(m.group(1))
            continue

        # Skip indicator
        stripped = line.strip()
        if stripped in ("SKIP", "**NO TRANSLATION NEEDED**", "NO TRANSLATION NEEDED"):
            _save()
            current_slide = None
            continue

        if current_slide is None:
            continue

        # Shape header (with or without ### prefix)
        m = re.match(r"^(?:###\s+)?Shape:\s*(.+)", line)
        if m:
            _save()
            current_shape = m.group(1).strip()
            continue

        # Paragraph — bullet format
        m = re.match(r"^-\s?(.*)", line)
        if m:
            spec = _parse_bold_text(m.group(1))
            current_paras.append(spec)
            continue

        # Paragraph — "Text: ..." format
        m = re.match(r"^Text:\s*(.*)", line)
        if m:
            spec = _parse_bold_text(m.group(1).strip())
            current_paras.append(spec)
            continue

        # Template lines (> ...) — skip (these are originals, not translations)
        if line.startswith(">"):
            continue

    _save()
    return result


# ─── PPTX manipulation helpers ────────────────────────────────

def _clone_rpr(source_rPr, src_lang="he-IL", tgt_lang="en-US"):
    """Deep-copy a <a:rPr>, replacing source language with target."""
    if source_rPr is None:
        rPr = etree.SubElement(etree.Element("dummy"), qn("a:rPr"))
        rPr.set("lang", tgt_lang)
        rPr.set("dirty", "0")
        return rPr
    new = copy.deepcopy(source_rPr)
    if new.get("lang", "") == src_lang:
        new.set("lang", tgt_lang)
    if new.get("altLang", "") == src_lang:
        new.set("altLang", tgt_lang)
    return new


def _clear_runs(p_elem):
    """Remove all <a:r>, <a:br>, and <a:endParaRPr> from a paragraph."""
    for tag in (qn("a:r"), qn("a:br"), qn("a:endParaRPr")):
        for child in p_elem.findall(tag):
            p_elem.remove(child)


def _fix_ppr_ltr(p_elem):
    """Remove RTL attributes from paragraph properties."""
    pPr = p_elem.find(qn("a:pPr"))
    if pPr is not None:
        for attr in list(pPr.attrib.keys()):
            if "rtl" in attr.lower():
                del pPr.attrib[attr]


def _add_run(p_elem, text, template_rPr, overrides=None,
             src_lang="he-IL", tgt_lang="en-US"):
    """Append an <a:r> element to a paragraph."""
    r_elem = etree.SubElement(p_elem, qn("a:r"))

    new_rPr = _clone_rpr(template_rPr, src_lang, tgt_lang)
    new_rPr.set("lang", tgt_lang)
    new_rPr.set("dirty", "0")

    if overrides:
        if "bold" in overrides:
            if overrides["bold"]:
                new_rPr.set("b", "1")
            elif "b" in new_rPr.attrib:
                del new_rPr.attrib["b"]
        if overrides.get("color"):
            for sf in new_rPr.findall(qn("a:solidFill")):
                new_rPr.remove(sf)
            sf = etree.SubElement(new_rPr, qn("a:solidFill"))
            clr = etree.SubElement(sf, qn("a:srgbClr"))
            clr.set("val", overrides["color"])

    r_elem.insert(0, new_rPr)

    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = text
    if text and (text[0] == " " or text[-1] == " " or "  " in text):
        t_elem.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")


def _ensure_endpararpr(txBody, src_lang, tgt_lang):
    """Ensure <a:endParaRPr> is the LAST child of every <a:p>."""
    for p_elem in txBody.findall(qn("a:p")):
        epr = p_elem.find(qn("a:endParaRPr"))
        if epr is not None:
            p_elem.remove(epr)
            if epr.get("lang", "") == src_lang:
                epr.set("lang", tgt_lang)
            p_elem.append(epr)
        else:
            epr = etree.SubElement(p_elem, qn("a:endParaRPr"))
            epr.set("lang", tgt_lang)
            epr.set("dirty", "0")


# ─── Core shape translation ──────────────────────────────────

def _replace_shape_text(shape, para_specs, src_lang="he-IL", tgt_lang="en-US"):
    """Replace text in a shape while preserving original run formatting."""
    if not shape.has_text_frame:
        return

    txBody = shape.text_frame._txBody
    orig_p = txBody.findall(qn("a:p"))

    # Collect template rPr from each original paragraph's first run
    templates = []
    for p in orig_p:
        rPr = None
        for r in p.findall(qn("a:r")):
            rPr = r.find(qn("a:rPr"))
            if rPr is not None:
                break
        templates.append(rPr)

    target_count = len(para_specs)
    current_count = len(orig_p)

    # Add paragraphs if needed
    while current_count < target_count:
        new_p = copy.deepcopy(orig_p[-1]) if orig_p else etree.SubElement(txBody, qn("a:p"))
        _clear_runs(new_p)
        txBody.append(new_p)
        templates.append(templates[-1] if templates else None)
        orig_p = txBody.findall(qn("a:p"))
        current_count += 1

    # Remove excess paragraphs
    orig_p = txBody.findall(qn("a:p"))
    while current_count > target_count and target_count > 0:
        txBody.remove(orig_p[current_count - 1])
        orig_p = txBody.findall(qn("a:p"))
        templates.pop()
        current_count -= 1

    # Process each paragraph
    source_is_rtl = _is_rtl_lang(src_lang)
    orig_p = txBody.findall(qn("a:p"))

    for i, spec in enumerate(para_specs):
        p_elem = orig_p[i]
        tpl = templates[i] if i < len(templates) else (templates[-1] if templates else None)

        if source_is_rtl:
            _fix_ppr_ltr(p_elem)

        if spec is None:
            continue

        _clear_runs(p_elem)

        if spec == "":
            # Empty / spacer paragraph
            continue

        if isinstance(spec, str):
            _add_run(p_elem, spec, tpl, src_lang=src_lang, tgt_lang=tgt_lang)

        elif isinstance(spec, list):
            for run_spec in spec:
                if isinstance(run_spec, tuple) and len(run_spec) == 2:
                    text, overrides = run_spec
                    _add_run(p_elem, text, tpl, overrides,
                             src_lang=src_lang, tgt_lang=tgt_lang)
                elif isinstance(run_spec, str):
                    _add_run(p_elem, run_spec, tpl,
                             src_lang=src_lang, tgt_lang=tgt_lang)

    # Flip RTL RIGHT alignment → LTR LEFT
    if source_is_rtl:
        for p_elem in txBody.findall(qn("a:p")):
            pPr = p_elem.find(qn("a:pPr"))
            if pPr is not None and pPr.get("algn", "") == "r":
                pPr.set("algn", "l")

    _ensure_endpararpr(txBody, src_lang, tgt_lang)


def _fix_all_rtl(prs, src_lang="he-IL", tgt_lang="en-US"):
    """Sweep the entire presentation: remove RTL, fix lang tags."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                pPr = para._pPr
                if pPr is not None:
                    for attr in list(pPr.attrib.keys()):
                        if "rtl" in attr.lower():
                            del pPr.attrib[attr]
                for run in para.runs:
                    rPr = run._r.find(qn("a:rPr"))
                    if rPr is not None:
                        if rPr.get("lang", "") == src_lang:
                            rPr.set("lang", tgt_lang)
                        if rPr.get("altLang", "") == src_lang:
                            rPr.set("altLang", tgt_lang)


# ─── High-level API ───────────────────────────────────────────

def translate_pptx(input_path, md_text, output_path,
                   src_lang_code=None, tgt_lang_code="en-US"):
    """Translate a PPTX file using a translated MD, write to output_path."""
    prs = Presentation(input_path)
    if src_lang_code is None:
        src_lang_code = _detect_source_lang(prs) or "he-IL"

    translations = parse_translation_md(md_text)
    warnings = _apply(prs, translations, src_lang_code, tgt_lang_code)
    prs.save(output_path)
    return warnings


def translate_pptx_bytes(pptx_bytes, md_text,
                         src_lang_code=None, tgt_lang_code="en-US"):
    """Translate PPTX from bytes, return (output_bytes, warnings)."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    if src_lang_code is None:
        src_lang_code = _detect_source_lang(prs) or "he-IL"

    translations = parse_translation_md(md_text)
    warnings = _apply(prs, translations, src_lang_code, tgt_lang_code)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue(), warnings


def _apply(prs, translations, src_lang_code, tgt_lang_code):
    """Apply parsed translations to a Presentation object."""
    warnings = []

    for slide_num, shape_list in sorted(translations.items()):
        slide_idx = slide_num - 1
        if slide_idx < 0 or slide_idx >= len(prs.slides):
            warnings.append(f"Slide {slide_num} does not exist in the PPTX.")
            continue

        slide = prs.slides[slide_idx]

        for shape_name, para_specs in shape_list:
            target = None
            for shape in slide.shapes:
                if shape.name == shape_name:
                    target = shape
                    break
            if target is None:
                warnings.append(
                    f"Slide {slide_num}: shape '{shape_name}' not found."
                )
                continue

            _replace_shape_text(target, para_specs, src_lang_code, tgt_lang_code)

    if _is_rtl_lang(src_lang_code):
        _fix_all_rtl(prs, src_lang_code, tgt_lang_code)

    return warnings


# ─── ChatGPT prompt ───────────────────────────────────────────

def get_chatgpt_prompt(source_lang="Hebrew", target_lang="English"):
    return f"""I have a presentation in {source_lang} that needs to be translated to {target_lang}.

I'm providing:
1. The **PDF** of the presentation (attached file) — use this to read the actual text.
2. A **translation template** (attached file or pasted below) — extracted from the PowerPoint file with exact shape names and paragraph structure.

**Please translate the template following these rules:**

1. **Keep all headers exactly as-is**: `## Slide N` and `### Shape: <name>` must NOT be changed.
2. **Replace each `>` line** with a `- ` line containing the {target_lang} translation.
3. **Keep `SKIP`** for slides that do not need translation.
4. **Translate the FULL content** of every paragraph. Do NOT abbreviate, summarize, or shorten — keep ALL details from the PDF.
5. **Empty paragraphs**: for blank `>` lines, write just `-`
6. **Bold formatting**: use double asterisks — `- **This text is bold**`
7. **Mixed formatting**: `- Normal text **bold part** more normal text`
8. **Keep text already in {target_lang}** unchanged.
9. **Line breaks within one paragraph**: use `\\n` (literal backslash-n) in the text.
10. **Same paragraph count**: each shape must have the EXACT same number of `-` lines as `>` lines in the template.

**IMPORTANT**: Provide the COMPLETE, DETAILED translation. Do not shorten or omit any content. Use the PDF to read the full original text for each shape.

**Output**: Return ONLY the filled template with no extra commentary."""
