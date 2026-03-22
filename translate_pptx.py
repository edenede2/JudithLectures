"""
translate_pptx.py — Translate a Hebrew PPTX to English.

Usage:
    python translate_pptx.py <input.pptx> <translations.md> <output.pptx>

The translations.md file should follow the format documented at the bottom of this script.
The script copies the original PPTX (preserving all images, shapes, layouts, etc.)
and replaces only the text content with the English translations, adjusting
text direction from RTL to LTR where needed.
"""

import sys
import copy
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.oxml.ns import qn
from lxml import etree


# ──────────────────────────────────────────────────────────────────────
# Translation map: slide_number -> list of shape translations
# Each shape translation is a dict with:
#   "shape_name" or "shape_idx": identifier
#   "paragraphs": list of paragraph dicts
# Each paragraph dict has:
#   "runs": list of run dicts, each with "text", optionally "bold", "size", "color", "lang"
#   Optionally "align": "LEFT", "CENTER", "RIGHT"
# ──────────────────────────────────────────────────────────────────────

def build_translation_map():
    """
    Build a slide-by-slide translation map.
    Returns dict: slide_number (1-based) -> list of (shape_match, paragraphs) tuples.
    
    shape_match can be:
      {"name": "Title 1"} — match by shape name
      {"idx": 0} — match by shape index
      
    paragraphs is a list of paragraph specs. Each paragraph spec:
      (alignment, [(text, bold, font_size_emu, color_hex, lang), ...])
    
    alignment: "LEFT", "CENTER", "RIGHT", None (preserve original)
    font_size_emu: None means preserve original
    color_hex: None means preserve original, e.g. "FF0000"
    lang: language tag, e.g. "en-US"
    """
    T = {}
    
    # Helper to make run tuples simpler
    def r(text, bold=None, size=None, color=None, lang="en-US"):
        return (text, bold, size, color, lang)

    # ── SLIDE 1 ──
    T[1] = [
        ({"name": "Title 1"}, [
            ("CENTER", [r("Project Life Cycle in Data Science")]),
        ]),
        ({"name": "Subtitle 2"}, [
            ("CENTER", [r("Judith Somekh", size=647700)]),
        ]),
    ]
    
    # ── SLIDE 2 ── (title already in English, content already in English)
    # Title and bullets are already English in original. Skip.
    
    # ── SLIDE 3 ── (diagram only, no text to translate)
    
    # ── SLIDE 4 ── section divider
    T[4] = [
        ({"name": "Title 1"}, [
            ("LEFT", [
                r("2. Data Gathering & Understanding", size=685800),
                r("\n", size=685800),
                r("Data collection & understanding", size=685800),
            ]),
        ]),
    ]
    
    # ── SLIDE 5 ──
    T[5] = [
        ({"name": "Content Placeholder 2"}, [
            ("LEFT", [r("The data include the raw material available from which the solution will be built.")]),
            ("LEFT", [r("Usually, finding the right data takes time and effort.")]),
            ("LEFT", [r("Data are often collected from multiple sources.")]),
            None,  # empty paragraph preserved
            ("LEFT", [r("What are some of the problems that can exist in data?")]),
        ]),
    ]
    
    # ── SLIDE 6 ──
    T[6] = [
        ({"name": "Content Placeholder 2"}, [
            ("LEFT", [r("Some of the questions worth considering are:", size=457200)]),
            ("LEFT", [r("What data do I need for my project?", size=457200)]),
            ("LEFT", [r("Where are the data located?", size=457200)]),
            ("LEFT", [r("How can I obtain them?", size=457200)]),
            ("LEFT", [r("What are the costs?", size=457200)]),
            ("LEFT", [r("What is the most efficient way to store and access all of this?", size=457200)]),
            ("LEFT", [r("How do we preserve data privacy and confidentiality?", size=457200)]),
        ]),
    ]
    
    # ── SLIDE 7 ──
    T[7] = [
        ({"name": "Content Placeholder 2"}, [
            ("LEFT", [r("It is important to understand the strengths and limitations of the data, because there is rarely an exact match between the data and the problem.")]),
            ("LEFT", [r("Historical data are often collected for purposes unrelated to the current business problem, or with no explicit purpose at all (for example, medical records).", size=254000)]),
            ("LEFT", [r("For example, a customer database, a transaction database, and a marketing response database contain different information, may cover different overlapping populations, and may have different levels of reliability.", size=254000)]),
            ("LEFT", [r("Data costs may vary. Some data will be available for free, while others will need to be purchased or will require effort to obtain (for example, GEO, the Israeli Digital Health Project).")]),
            ("LEFT", [r("Sometimes the data do not exist, and entire projects will be required to organize their collection (for example, gene expression data require performing experiments).", size=254000)]),
            ("LEFT", [r("A critical part of the data understanding stage is estimating the costs and benefits of each data source.", size=254000)]),
        ]),
    ]
    
    # ── SLIDE 8 ──
    T[8] = [
        ({"name": "Title 1"}, [
            ("LEFT", [r("Free medical information for researchers and start-up companies")]),
        ]),
    ]
    
    # ── SLIDE 9 ──
    T[9] = [
        ({"name": "Content Placeholder 2"}, [
            ("LEFT", [r("As data understanding progresses, solution paths may change direction in response, and the team's efforts may even split.")]),
            ("LEFT", [r("Example: Fraud detection \u2013 data mining is widely used for fraud detection, and many fraud detection problems include classic data mining tasks.")]),
            ("LEFT", [r("Sometimes fraud problems that initially look similar turn out to be fundamentally different during the data understanding process.")]),
            ("LEFT", [r("Catching credit card fraud:", bold=True)]),
            ("LEFT", [r("How would we do that?")]),
        ]),
    ]
    
    # ── SLIDE 10 ──
    T[10] = [
        ({"name": "Content Placeholder 2"}, [
            ("LEFT", [r("Catching credit card fraud:")]),
            ("LEFT", [r("To build a predictive model for fraud, we need to label actions as fraudulent / legitimate.")]),
            ("LEFT", [r("Question: Who will perform the labeling?")]),
            ("LEFT", [r("How will we identify fraud? Who will identify it?")]),
            ("LEFT", [r("Is the identification reliable?")]),
        ]),
    ]
    
    # ── SLIDE 11 ──
    T[11] = [
        ({"name": "Content Placeholder 2"}, [
            ("LEFT", [r("Catching credit card fraud:")]),
            ("LEFT", [r("Charges appear in each customer\u2019s account, so fraudulent charges are usually detected\u2014if not initially by the company, then later by the customer when reviewing account activity.")]),
            ("LEFT", [r("We can assume that almost all fraud is detected and labeled reliably, because the legitimate customer and the person committing the fraud are different people with opposite interests.")]),
            ("LEFT", [r("Credit card transactions have reliable labels (fraud or legitimate) that can be used to train a machine learning model.", size=355600)]),
        ]),
    ]
    
    # ── SLIDE 12 ──
    T[12] = [
        ({"name": "Content Placeholder 2"}, [
            None, None,  # empty paragraphs
            ("LEFT", [r("Medicare fraud: this is a huge problem in the United States, costing billions of dollars each year.", size=406400)]),
            ("LEFT", [r("The fraudsters are healthcare providers who submit false claims for various treatments for patients.", size=406400)]),
            None,  # empty paragraph
            ("LEFT", [r("Is this a conventional fraud detection problem?", size=406400)]),
        ]),
    ]
    
    # ── SLIDE 13 ──
    T[13] = [
        ({"name": "Content Placeholder 2"}, [
            None, None,
            ("LEFT", [r("Medicare fraud: this is a huge problem in the United States, costing billions of dollars each year.", size=406400)]),
            ("LEFT", [r("The fraudsters are healthcare providers who submit false claims for various treatments for patients.", size=406400)]),
            None,
            ("LEFT", [r("Is this a conventional fraud detection problem?", size=406400)]),
            ("LEFT", [r("When we examine the relationship between the business problem and the data, we understand that the problem is significantly different.", size=304800)]),
            ("LEFT", [r("Why? How will we identify fraud? Who identifies it?", size=406400)]),
        ]),
    ]
    
    # ── SLIDE 14 ──
    T[14] = [
        ({"name": "Content Placeholder 2"}, [
            None, None,
            ("LEFT", [r("Example: Medicare fraud")]),
            ("LEFT", [r("The perpetrators of fraud\u2014healthcare providers submitting false claims, and sometimes their patients\u2014are also legitimate service providers using the billing system.", size=304800)]),
            ("LEFT", [
                r("The people committing the fraud are a subset of legitimate users; there is no separate disinterested party that will accurately declare which claims are \u201ccorrect.\u201d", bold=True, size=304800),
            ]),
            ("LEFT", [
                r("As a result, Medicare billing data do not contain a reliable target variable indicating fraud, so a supervised learning approach that might work for credit card fraud is not relevant here.", size=304800),
            ]),
            ("LEFT", [
                r("Such a problem usually requires unsupervised approaches such as clustering, anomaly detection, and co-occurrence grouping of similar data.", size=304800),
            ]),
        ]),
    ]
    
    # ── SLIDE 15 ──
    T[15] = [
        ({"name": "Content Placeholder 2"}, [
            ("LEFT", [r("Statistical learning lies at the intersection of computer science and statistics and refers to a set of tools for understanding data. It combines statistical methods, computational methods, and methods from machine learning.")]),
            ("LEFT", [r("These tools are categorized as:")]),
            ("LEFT", [r("Supervised learning", bold=True)]),
            ("LEFT", [r("Building a statistical model for predicting or estimating an output based on one or more inputs.")]),
            ("LEFT", [r("Unsupervised learning", bold=True)]),
            ("LEFT", [r("Inputs exist, but there is no known corresponding output. Even so, we can learn about relationships and structure from the data. There is no response variable that \u201csupervises\u201d the analysis.")]),
        ]),
    ]
    
    # ── SLIDE 16 ──
    T[16] = [
        ({"name": "Content Placeholder 2"}, [
            ("LEFT", [
                r("The fact that both cases are called fraud detection problems is based on a "),
                r("superficial similarity that is actually misleading", bold=True),
                r(", because they are two different cases requiring "),
                r("different analyses", bold=True),
                r("."),
            ]),
            ("LEFT", [
                r("In data understanding, ", bold=True),
                r("we must dig beneath the surface to uncover the structure of the business problem and the available data for understanding it, and then match them to one or more data mining and analysis tasks requiring the knowledge, science, and technology we possess."),
            ]),
            ("LEFT", [r("It is not unusual for a business problem to contain several data mining and analysis tasks, often of different kinds, requiring integration of the solutions to the different tasks.")]),
        ]),
    ]
    
    # ── SLIDE 17 ── (diagram only)
    
    # ── SLIDE 18 ── section divider
    T[18] = [
        ({"name": "Title 1"}, [
            ("LEFT", [
                r("3. Data Cleaning & Preparation", size=609600),
                r("\n", size=609600),
                r("Cleaning and preparing data", size=609600),
            ]),
        ]),
    ]
    
    # ── SLIDE 19 ──
    T[19] = [
        ({"name": "Content Placeholder 2"}, [
            ("LEFT", [r("Powerful data analysis techniques exist, but they are still limited by certain requirements imposed on the data they use.", size=406400)]),
            ("LEFT", [r("Data entered into analysis usually need to be represented differently from the original form in which they were collected or delivered, and this often requires conversion and cleaning efforts.", size=406400)]),
            ("LEFT", [r("What kinds of problems can our data have?", size=406400)]),
        ]),
    ]
    
    # ── SLIDE 20 ── (mostly English already, skip)
    
    # ── SLIDE 21 ── (mostly English already, skip)
    
    # ── SLIDE 22 ──
    T[22] = [
        ({"name": "Title 1"}, [
            ("CENTER", [r("A Mars lander crashed because of a unit conversion error", color="FFFFFF")]),
        ]),
    ]
    
    # ── SLIDE 23 ──
    T[23] = [
        ({"name": "Title 1"}, [
            ("LEFT", [r("Data Leakage", size=685800)]),
        ]),
        ({"name": "Content Placeholder 2"}, [
            ("LEFT", [r("Leakage is a situation in which a variable collected in historical data provides information about the target variable\u2014information that appears in the historical data but is not actually available when a decision must be made.", size=304800)]),
            ("LEFT", [r("Example 1: We want to predict whether a customer will be a \u201cbig spender\u201d; knowing the categories of items purchased (or the amount of tax paid) may be highly predictive, but these are not known at decision time (Kohavi & Parekh, 2003).", size=304800)]),
            ("LEFT", [r("Example 2: When predicting whether, at a certain point in time, a website visitor will end her session or continue to another page, the variable \u201ctotal number of webpages visited in the session\u201d is predictive. However, that total is only known after the session ends (Kohavi et al., 2000)\u2014at which point the target is already known.", size=304800)]),
        ]),
    ]
    
    # ── SLIDE 24 ── (diagram only)
    
    # ── SLIDE 25 ── section divider
    T[25] = [
        ({"name": "Title 1"}, [
            ("LEFT", [
                r("4. Data Exploration & Modeling", bold=True, size=609600),
                r("\n", size=609600),
                r("Data exploration and models", size=609600),
            ]),
        ]),
    ]
    
    # ── SLIDE 26 ──
    T[26] = [
        ({"name": "Content Placeholder 3"}, [
            ("LEFT", [r("The actual data analysis process is based on computational techniques, statistical methods, and algorithms.", size=508000)]),
            ("LEFT", [r("The result of modeling is some form of model that finds patterns and structures in the data.", size=508000)]),
        ]),
    ]
    
    # ── SLIDE 27 ──
    T[27] = [
        ({"name": "Content Placeholder 3"}, [
            ("LEFT", [
                r("We must "),
                r("understand the patterns, problems, and biases in the data", bold=True, size=508000),
                r(":", size=508000),
            ]),
            ("LEFT", [r("This is done by sampling and analyzing a random subset of the data from the full dataset.", size=457200)]),
            ("LEFT", [r("Use histograms or distribution curves to see the general trend, or even build an interactive visualization that allows drilling down to individual data points.", size=457200)]),
            ("LEFT", [r("Investigate the story behind outliers.", size=457200)]),
            ("LEFT", [
                r("Begin forming hypotheses", bold=True, size=508000),
                r(" about the data and the problem we are trying to solve or analyze.", size=508000),
            ]),
            ("LEFT", [r("Example: In predicting students\u2019 grades, you could examine the relationship between grades and sleep.", size=457200)]),
            ("LEFT", [r("Example: In predicting real-estate prices, you could draw prices as a spatial heatmap to see whether you can detect trends.", size=457200)]),
        ]),
    ]
    
    # ── SLIDE 28 ── (mostly English already, skip)
    
    # ── SLIDE 29 ──
    T[29] = [
        ({"name": "Content Placeholder 2"}, [
            ("LEFT", [r("Patients with cancer and healthy individuals")]),
            ("LEFT", [r("The features of each patient are represented by thousands of genes (measured mRNA levels), and the levels change between healthy and diseased states")]),
        ]),
    ]
    
    # ── SLIDE 30 ──
    T[30] = [
        ({"name": "Title 1"}, [
            ("LEFT", [
                r("4.2 Feature Engineering", bold=True),
            ]),
        ]),
        ({"name": "Content Placeholder 3"}, [
            ("LEFT", [r("In machine learning, a feature is a measurable property or characteristic of an observed phenomenon.")]),
            ("LEFT", [r("For predicting student grades, one possible feature is the amount of sleep at night.")]),
            ("LEFT", [r("Question: what other relevant feature ideas can you think of?")]),
            ("LEFT", [r("In more complex tasks such as face recognition, features can be histograms counting the number of black pixels.")]),
            ("LEFT", [r("According to Andrew Ng: \u201cComing up with features is difficult, time-consuming, requires expert knowledge. \u2018Applied machine learning\u2019 is basically feature engineering.\u201d")]),
        ]),
    ]
    
    # ── SLIDE 31 ──
    T[31] = [
        ({"name": "Title 1"}, [
            ("LEFT", [r("4.2 Feature Engineering", bold=True)]),
        ]),
        ({"name": "Content Placeholder 3"}, [
            ("LEFT", [r("Feature engineering is the process of using domain knowledge to transform raw data into informative features that represent the business problem being solved.")]),
            ("LEFT", [r("This stage directly affects the accuracy of the predictive model built in the next stage.")]),
            ("LEFT", [r("Usually we perform two types of tasks in feature engineering:")]),
            ("LEFT", [r("Feature selection \u2013 choosing features relevant to the analysis", bold=True)]),
            ("LEFT", [r("Feature construction \u2013 for example, combining features into one feature", bold=True)]),
        ]),
    ]
    
    # ── SLIDE 32 ──
    T[32] = [
        ({"name": "Title 1"}, [
            ("LEFT", [r("4.2 Feature Selection", bold=True)]),
        ]),
        ({"name": "Content Placeholder 3"}, [
            ("LEFT", [r("Feature selection is the process of removing features that add more noise than information.", size=304800)]),
            ("LEFT", [r("For this purpose we use filtering methods:", size=304800)]),
            ("LEFT", [r("Filter methods \u2013 generate a statistical measure to assign a score to each feature.", size=304800)]),
            ("LEFT", [
                r("Wrapper methods", size=304800),
                r(" \u2013 treat feature selection as a search problem and use heuristics to perform the search.", size=304800),
            ]),
            ("LEFT", [
                r("Embedded methods", size=304800),
                r(" \u2013 use machine learning itself to determine which features contribute best to accuracy.", size=304800),
            ]),
        ]),
    ]
    
    # ── SLIDE 33 ──
    T[33] = [
        ({"name": "Title 1"}, [
            ("LEFT", [r("4.2 Feature Construction", bold=True)]),
        ]),
        ({"name": "Content Placeholder 3"}, [
            ("LEFT", [r("Creating new features from original features.")]),
            ("LEFT", [r("For example, when you have a continuous variable but domain knowledge says that only crossing a certain threshold matters.")]),
            ("LEFT", [r("Given feature: subject\u2019s age.")]),
            ("LEFT", [r("If the model only cares whether a person is an adult or a minor, we can place a threshold at age 18 and assign different categories to cases above and below it.")]),
            ("LEFT", [r("Other threshold examples: sick/healthy, obese/thin.")]),
            ("LEFT", [r("Merging multiple features to make them more informative by taking their sum, difference, or product.")]),
            ("LEFT", [r("Example: predicting student grades.")]),
            ("LEFT", [r("Given: features representing number of hours of sleep for a student each night.")]),
            ("LEFT", [r("What processing would you perform on these features for prediction?")]),
        ]),
    ]
    
    # ── SLIDE 34 ──
    T[34] = [
        ({"name": "Title 1"}, [
            ("LEFT", [r("4.2 Feature Construction", bold=True)]),
        ]),
        ({"name": "Content Placeholder 3"}, [
            ("LEFT", [r("Creating new features from original features.")]),
            ("LEFT", [r("For example, when you have a continuous variable but domain knowledge says that only crossing a certain threshold matters.")]),
            ("LEFT", [r("Given feature: subject\u2019s age.")]),
            ("LEFT", [r("If the model only cares whether a person is an adult or a minor, we can place a threshold at age 18 and assign different categories to cases above and below it.")]),
            ("LEFT", [r("Other threshold examples: sick/healthy, obese/thin.")]),
            ("LEFT", [r("Merging multiple features to make them more informative by taking their sum, difference, or product.")]),
            ("LEFT", [r("Example: predicting student grades.")]),
            ("LEFT", [r("Given: features representing number of hours of sleep for a student each night.")]),
            ("LEFT", [r("What processing would you perform on these features for prediction?")]),
            ("LEFT", [r("You should create a feature representing the student\u2019s average amount of sleep.")]),
        ]),
    ]
    
    # ── SLIDE 35 ──
    T[35] = [
        ({"name": "Title 1"}, [
            ("LEFT", [r("4.3 Predictive Modeling", bold=True)]),
        ]),
        ({"name": "Content Placeholder 3"}, [
            ("LEFT", [r("The actual analysis work!")]),
            ("LEFT", [r("Includes:")]),
            ("LEFT", [r("Machine learning techniques for building a predictive model and testing its accuracy.")]),
            ("LEFT", [r("Use of comprehensive statistical methods and tests to ensure that the model\u2019s results are sensible and meaningful.")]),
            ("LEFT", [
                r("Based on the questions asked during the business understanding stage, this is where we decide "),
                r("which model to choose for the problem.", bold=True),
            ]),
            ("LEFT", [r("This is never an easy decision, and there is no single correct answer.")]),
            ("LEFT", [r("The chosen model (or models\u2014and it is always preferable to test several) will depend on the size, type, and quality of your data, how much time and computing resources you are willing to invest, and the type of output you intend to produce.")]),
        ]),
    ]
    
    # ── SLIDE 36 ──
    T[36] = [
        ({"name": "Content Placeholder 3"}, [
            ("LEFT", [r("After training your model, it is critical to evaluate its quality and its success in prediction.")]),
            ("LEFT", [r("Cross-validation\u2014usually k-fold cross-validation\u2014is commonly used to measure model accuracy.")]),
            ("LEFT", [r("This involves splitting the dataset into equally sized groups of instances, building a model on all but one group, and repeating the process with different held-out groups. This allows the model to be trained across all the data rather than relying on a simple train-test split.")]),
            ("LEFT", [r("Charts such as ROC curves, which show the true positive rate against the false positive rate, are also used to measure model success.")]),
        ]),
    ]
    
    # ── SLIDE 37 ──
    T[37] = [
        ({"name": "Title 1"}, [
            ("LEFT", [
                r("4.4 ", bold=True),
                r("Visualization"),
            ]),
        ]),
        ({"name": "Content Placeholder 3"}, [
            ("LEFT", [r("Data visualization seems simple, but often it is one of the most challenging tasks.")]),
            ("LEFT", [r("Data visualization combines communication, psychology, statistics, and art.")]),
            ("LEFT", [r("The ultimate goal is to communicate the data in a simple, effective, and visually pleasing way.")]),
            ("LEFT", [r("After extracting the intended insights from your model, you must represent them in a way that the various stakeholders in the project can understand.")]),
            ("LEFT", [r("The chart should be clear and tell the story!")]),
        ]),
    ]
    
    # ── SLIDE 38 ── (diagram only)
    
    # ── SLIDE 39 ──
    T[39] = [
        ({"name": "TextBox 6"}, [
            ("LEFT", [r("Goal 1: Rigorously evaluate the results of the data analysis and ensure that they are valid and reliable before moving on.", size=355600)]),
            ("LEFT", [r("If we search enough, we will find patterns in any dataset\u2014but they may not survive more careful testing.", size=355600)]),
            ("LEFT", [
                r("We want to be sure that the models and patterns extracted from the data are ", size=355600),
                r("real", size=355600),
                r(" regularities and not just strange phenomena or sampling artifacts.", size=355600),
            ]),
        ]),
    ]
    
    # ── SLIDE 40 ──
    T[40] = [
        ({"name": "TextBox 6"}, [
            None,  # empty P0
            ("LEFT", [
                r("Controlled evaluation", bold=True, size=304800),
                r(": first test a model within a controlled \u201claboratory\u201d process.", size=304800),
            ]),
            ("LEFT", [
                r("Evaluation of external considerations", bold=True, size=304800),
                r(": even if a model passes strict evaluation tests in a controlled environment, there may still be external considerations that make it impractical.", size=304800),
            ]),
            ("LEFT", [r("For example, a common flaw in detection analyses (such as fraud detection, spam detection, and intrusion monitoring) is that they generate too many false alarms.", size=304800)]),
            ("LEFT", [r("Additional considerations must be addressed, such as: how much it will cost the team to deal with all those false alarms, and what the meaning and cost of customer dissatisfaction will be.", size=304800)]),
        ]),
    ]
    
    # ── SLIDE 41 ──
    T[41] = [
        ({"name": "TextBox 6"}, [
            ("LEFT", [
                r("Evaluation in the real environment", bold=True, size=304800),
                r(": in some cases we will extend evaluation into the real environment, for example by instrumenting a live system in the real setting to allow randomized experiments.", size=304800),
            ]),
            None, None,  # empty paragraphs
            ("LEFT", [r("Why do you think this is necessary?", size=304800)]),
        ]),
    ]
    
    # ── SLIDE 42 ──
    T[42] = [
        ({"name": "TextBox 6"}, [
            ("LEFT", [
                r("Evaluation in the real environment", bold=True, size=304800),
                r(": in some cases we will extend evaluation into the real environment, for example by using instrumentation in a live system that enables randomized experiments.", size=304800),
            ]),
            ("LEFT", [r("In the example of customer churn in cellular companies\u2014if laboratory testing suggests that a data-driven model can help reduce churn, we may want to move to real-world evaluation in a live system that applies the model randomly to some customers while keeping others as a control group. Such experiments must be carefully planned.", size=304800)]),
            ("LEFT", [
                r("Considering the changing environment", bold=True, size=304800),
                r(":", size=304800),
            ]),
            ("LEFT", [r("What might change when the model is implemented in the company?", size=304800)]),
        ]),
    ]
    
    # ── SLIDE 43 ──
    T[43] = [
        ({"name": "TextBox 6"}, [
            ("LEFT", [
                r("Evaluation in the real environment", bold=True, size=304800),
                r(": in some cases we will extend evaluation into the real environment, for example by using instrumentation in a live system that enables randomized experiments.", size=304800),
            ]),
            ("LEFT", [r("In the example of customer churn in cellular companies\u2014if laboratory testing suggests that a data-driven model can help reduce churn, we may want to move to real-world evaluation in a live system that applies the model randomly to some customers while keeping others as a control group. Such experiments must be carefully planned.", size=304800)]),
            ("LEFT", [
                r("Considering the changing environment", bold=True, size=304800),
                r(": We may also want to train our systems to ensure that the world does not change in a way that harms the model\u2019s decision-making. For example, behavior can change\u2014in some cases, such as ", size=304800),
                r("fraud or spam", bold=True, size=304800),
                r(", directly in response to deployment of the models.", size=304800),
            ]),
        ]),
    ]
    
    # ── SLIDE 44 ──
    T[44] = [
        ({"name": "Content Placeholder 2"}, [
            ("LEFT", [r("We want to ensure that the model meets the original business goals.", bold=True, size=304800)]),
            ("LEFT", [r("After going through the entire lifecycle, it is time to return to the drawing board. Remember, this is a circular lifecycle, so it is an iterative process. This is the point where you evaluate how the success of your model relates back to your original understanding of the business and its goals.", size=304800)]),
            ("LEFT", [r("Does it address the problems that were identified?", size=304800)]),
            ("LEFT", [r("Does the analysis yield solid solutions?", size=304800)]),
            ("LEFT", [
                r("If new insights arose during the first iteration of the lifecycle (which usually happens), you can now ", size=304800),
                r("inject that new knowledge", bold=True, size=304800),
                r(" into the next iteration in order to produce even stronger insights and use the power of data to generate strong and exceptional outcomes for your business or project.", size=304800),
            ]),
        ]),
    ]
    
    # ── SLIDE 45 ──
    T[45] = [
        ({"name": "Title 1"}, [
            ("CENTER", [
                r("Q: What do data scientists spend the most time doing?", bold=True),
                r("\n"),
                r("What percentage of the time do you think each task takes?"),
            ]),
        ]),
    ]
    
    # ── SLIDE 46 ── (already English)
    # ── SLIDE 47 ── (already English)
    
    return T


def set_paragraph_alignment(para, align_str):
    """Set paragraph alignment and remove RTL/bidi settings."""
    from pptx.enum.text import PP_ALIGN
    align_map = {
        "LEFT": PP_ALIGN.LEFT,
        "CENTER": PP_ALIGN.CENTER,
        "RIGHT": PP_ALIGN.RIGHT,
    }
    if align_str and align_str in align_map:
        para.alignment = align_map[align_str]


def remove_rtl_from_paragraph(para):
    """Remove RTL/bidi attributes from a paragraph's XML."""
    pPr = para._pPr
    if pPr is not None:
        # Remove rtl attribute
        for attr in list(pPr.attrib.keys()):
            if 'rtl' in attr.lower():
                del pPr.attrib[attr]
        # Remove algn=r if we're changing to LTR
        # (handled by set_paragraph_alignment instead)


def clear_paragraph_runs(para):
    """Remove all runs from a paragraph, preserving the paragraph element."""
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    # Remove all <a:r> elements
    for r_elem in para._p.findall(qn('a:r')):
        para._p.remove(r_elem)
    # Remove all <a:br> elements
    for br_elem in para._p.findall(qn('a:br')):
        para._p.remove(br_elem)


def add_run_to_paragraph(para, text, bold=None, font_size=None, color_hex=None, lang="en-US"):
    """Add a new run to an existing paragraph with specified formatting."""
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    
    # Create <a:r> element
    r_elem = etree.SubElement(para._p, qn('a:r'))
    
    # Create <a:rPr> for run properties
    rPr = etree.SubElement(r_elem, qn('a:rPr'))
    rPr.set('lang', lang)
    rPr.set('dirty', '0')
    
    if bold is True:
        rPr.set('b', '1')
    elif bold is False:
        rPr.set('b', '0')
    
    if font_size is not None:
        # font size in hundredths of a point
        rPr.set('sz', str(int(font_size / 12.7)))  # EMU to hundredths of pt
    
    if color_hex:
        solidFill = etree.SubElement(rPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', color_hex)
    
    # Create <a:t> for text
    t_elem = etree.SubElement(r_elem, qn('a:t'))
    t_elem.text = text
    # Preserve spaces
    if text and (text[0] == ' ' or text[-1] == ' ' or '  ' in text):
        t_elem.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')


def replace_shape_text(shape, paragraphs_spec):
    """
    Replace all text in a shape's text frame with the translated text.
    
    paragraphs_spec: list of paragraph specs.
    Each spec is either None (skip/preserve empty) or:
      (alignment, [(text, bold, size, color, lang), ...])
    """
    if not shape.has_text_frame:
        return
    
    tf = shape.text_frame
    existing_paras = list(tf.paragraphs)
    
    # We need to match the paragraph count.
    # Strategy: clear existing paragraphs and set new content.
    # We'll work with the XML directly for more control.
    
    txBody = tf._txBody
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    
    # Get all existing <a:p> elements
    existing_p_elems = txBody.findall(qn('a:p'))
    
    # We might need more or fewer paragraphs
    target_count = len(paragraphs_spec)
    current_count = len(existing_p_elems)
    
    # Add more paragraph elements if needed
    while current_count < target_count:
        new_p = copy.deepcopy(existing_p_elems[-1]) if existing_p_elems else etree.SubElement(txBody, qn('a:p'))
        # Clear runs from cloned paragraph
        for r_elem in new_p.findall(qn('a:r')):
            new_p.remove(r_elem)
        for br_elem in new_p.findall(qn('a:br')):
            new_p.remove(br_elem)
        txBody.append(new_p)
        existing_p_elems = txBody.findall(qn('a:p'))
        current_count = len(existing_p_elems)
    
    # Remove extra paragraphs if we have too many
    while current_count > target_count and target_count > 0:
        txBody.remove(existing_p_elems[-1])
        existing_p_elems = txBody.findall(qn('a:p'))
        current_count = len(existing_p_elems)
    
    # Now process each paragraph
    existing_p_elems = txBody.findall(qn('a:p'))
    for i, spec in enumerate(paragraphs_spec):
        if i >= len(existing_p_elems):
            break
            
        p_elem = existing_p_elems[i]
        para = tf.paragraphs[i]
        
        if spec is None:
            # Preserve empty paragraph - just clear any text
            clear_paragraph_runs(para)
            remove_rtl_from_paragraph(para)
            continue
        
        alignment, runs = spec
        
        # Remove RTL settings
        remove_rtl_from_paragraph(para)
        
        # Set alignment
        set_paragraph_alignment(para, alignment)
        
        # Clear existing runs
        clear_paragraph_runs(para)
        
        # Add new runs
        for run_spec in runs:
            text, bold, size, color, lang = run_spec
            add_run_to_paragraph(para, text, bold=bold, font_size=size, color_hex=color, lang=lang)


def fix_all_rtl(prs):
    """
    Go through ALL shapes in all slides and convert any remaining RTL paragraphs to LTR.
    This catches shapes we didn't explicitly translate but that have RTL settings.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    pPr = para._pPr
                    if pPr is not None:
                        # Remove rtl attribute
                        for attr in list(pPr.attrib.keys()):
                            if 'rtl' in attr.lower():
                                del pPr.attrib[attr]
                    # Fix run language tags from he-IL to en-US
                    for run in para.runs:
                        rPr = run._r.find(qn('a:rPr'))
                        if rPr is not None:
                            lang = rPr.get('lang', '')
                            if lang == 'he-IL':
                                rPr.set('lang', 'en-US')


def translate_pptx(input_path, output_path):
    """Main function to translate the PPTX."""
    print(f"Loading: {input_path}")
    prs = Presentation(input_path)
    
    translation_map = build_translation_map()
    
    for slide_num, shape_translations in translation_map.items():
        slide_idx = slide_num - 1
        if slide_idx >= len(prs.slides):
            print(f"  WARNING: Slide {slide_num} does not exist, skipping.")
            continue
        
        slide = prs.slides[slide_idx]
        print(f"  Translating slide {slide_num}...")
        
        for shape_match, paragraphs_spec in shape_translations:
            # Find the shape
            target_shape = None
            if "name" in shape_match:
                for shape in slide.shapes:
                    if shape.name == shape_match["name"]:
                        target_shape = shape
                        break
                if target_shape is None:
                    print(f"    WARNING: Shape '{shape_match['name']}' not found on slide {slide_num}")
                    continue
            elif "idx" in shape_match:
                idx = shape_match["idx"]
                if idx < len(slide.shapes):
                    target_shape = list(slide.shapes)[idx]
                else:
                    print(f"    WARNING: Shape index {idx} out of range on slide {slide_num}")
                    continue
            
            replace_shape_text(target_shape, paragraphs_spec)
    
    # Fix any remaining RTL issues across entire presentation
    print("  Fixing remaining RTL settings...")
    fix_all_rtl(prs)
    
    print(f"Saving: {output_path}")
    prs.save(output_path)
    print("Done!")


if __name__ == "__main__":
    if len(sys.argv) == 3:
        input_path = sys.argv[1]
        output_path = sys.argv[2]
    elif len(sys.argv) == 1:
        # Default paths
        input_path = "Lecture2_intro.pptx"
        output_path = "Lecture2_intro_EN.pptx"
    else:
        print("Usage: python translate_pptx.py <input.pptx> <output.pptx>")
        sys.exit(1)
    
    translate_pptx(input_path, output_path)
